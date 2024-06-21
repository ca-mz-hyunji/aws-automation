import json
import boto3
import os
import io
import math
import logging
import copy
import time
from collections import Counter
from datetime import date, timedelta, datetime
from dateutil.relativedelta import relativedelta
from pptx import Presentation
from pptx.chart.data import CategoryChartData, ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.parts.chart import ChartPart
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches
from pptx.shapes.autoshape import Shape
from pptx.parts.embeddedpackage import EmbeddedXlsxPart
from pptx.dml.color import RGBColor

# Logging
logger = logging.getLogger()
logger.setLevel(logging.INFO)


# 에러난 위치 확인하는 거
# except Exception as e:
# logger.error('{}'.format(trackback.format_exc()))
class ConfigData:
    def event_format(event, region):
        
        '''
        # key 변환 및 this 기준으로 sort
        # 데이터 예시
        {
        	'ap-northeast-2': {
        		'cache_count': {
        			'Redis': {
        				'this': 3,
        				'last': 2,
        				'before_last': 1
        			}
        		},
        		'cache_type': {
        			'cache.m3.medium': {
        				'this': 10,
        				'last': 10,
        				'before_last': 10
        			
        		}
        	}
        }
        '''
        
        if region == None:
            pass
        
        # region data가 있는지 확인
        elif region == True:
            # region을 key로 데이터 재생성
            formed_data = {}
            since_items = ['month_data_0', 'month_data_1', 'month_data_2', 'month_data_3']
            for since_item in since_items:
                if since_item in event['data']:
                    for region in event['data'][since_item]:
                        if region not in formed_data:
                            formed_data[region] = {} 
                        for title in event['data'][since_item][region]:
                            if title not in formed_data[region]:
                                formed_data[region][title] = {}
                            for data in event['data'][since_item][region][title]:
                                if data in formed_data[region][title]:
                                    formed_data[region][title][data][since_item] = event['data'][since_item][region][title][data]
                                else:
                                    formed_data[region][title][data] = {since_item: event['data'][since_item][region][title][data]}
            
            # $since_items 없는 데이터 채우기
            for region in formed_data:
                for title in formed_data[region]:
                    for data in formed_data[region][title]:
                        for since_item in since_items:
                            if since_item not in formed_data[region][title][data]:
                                formed_data[region][title][data][since_item] = 0
            
            # sort formed_data
            for region in formed_data:
                for title in formed_data[region]:
                    formed_data[region][title] = dict(sorted(formed_data[region][title].items(), key=lambda x: x[1]['month_data_0'], reverse=True))
            return formed_data
        
        elif region == False:
            pass
            # region이 없는 데이터는 커스텀이 각자 필요하여 각각 생성하였습니다.
    
    def data_maxlen_format(formed_data, region, key, maxlen):
        # *** 함수를 사용하기 전 event_format 함수를 사용할 것 ***
        if region == True:
            # cache_type이 7개를 초과 할 경우 other로 분류
            since_items = ['month_data_0', 'month_data_1', 'month_data_2', 'month_data_3']
            if key == None:
                return formed_data
            else:
                for region in formed_data:
                    if key in formed_data[region]:
                        if len(formed_data[region][key]) > maxlen:
                            # other 데이터 생성
                            formed_data[region][key]['other'] = {'month_data_0': 0,'month_data_1': 0,'month_data_2': 0,'month_data_3': 0}
                            for other_cache_type in list(formed_data[region][key])[(maxlen-1):-1]:
                                for since_item in since_items:
                                    formed_data[region][key]['other'][since_item] += formed_data[region][key][other_cache_type][since_item]
                                formed_data[region][key].pop(other_cache_type)
                                
            return formed_data
        
        elif region == False:
            pass
    
    def date_format(run_date, report_cycle):
        if type(run_date) is not date:
            run_date = datetime.strptime(run_date, '%Y-%m-%d').date()
        
        run_date = run_date.replace(day=1)
        if report_cycle == '월간':
            info_date = str(run_date) + ' ~ ' + str(run_date + relativedelta(months=1) - timedelta(days=1))
        else:
            info_date = str(run_date - relativedelta(months=2)) + ' ~ ' + str(run_date + relativedelta(months=1) - timedelta(days=1))
        main_date = str(run_date.year) + '년 ' + str(run_date.month) + '월'
        
        file_date = main_date.replace(' ', '_')
        
        table_date = [
            (run_date-relativedelta(months=3)).strftime('%y.%m'),
            (run_date-relativedelta(months=2)).strftime('%y.%m'),
            (run_date-relativedelta(months=1)).strftime('%y.%m'),
            run_date.strftime('%y.%m'),
        ]
        return {'main': main_date, 'info': info_date, 'file': file_date, 'table': table_date}
    
    def crease(data, report_cycle):
        if report_cycle == '월간':
            compare_target = data.get('last', 0)    
        elif report_cycle == '분기':
            compare_target = data.get('before_last', 0)   
        
        result = data.get('this') - compare_target
        if result > 0:
            return '+'+str(result)
        elif result < 0:
            return str(result)
        else:
            return ''
    
    def add_event_data(event, slide_idx):
        agenda_list = []
        # 날짜 데이터 추가
        # 이번달
        if 'month_data_0' in event:
            for month_data_0 in event['month_data_0']:
                agenda_list.append(month_data_0)
                slide_idx[month_data_0]['data'] = {'month_data_0': event['month_data_0'][month_data_0]}
        
        # 1달전
        if 'month_data_1' in event:
            for month_data_1 in event['month_data_1']:
                agenda_list.append(month_data_1)
                if 'data' not in slide_idx[month_data_1]:
                    slide_idx[month_data_1]['data'] = {'month_data_1': event['month_data_1'][month_data_1]}
                else:
                    slide_idx[month_data_1]['data']['month_data_1'] = event['month_data_1'][month_data_1]
        
        # 2달전
        if 'month_data_2' in event:
            for month_data_2 in event['month_data_2']:
                agenda_list.append(month_data_2)
                if 'data' not in slide_idx[month_data_2]:
                    slide_idx[month_data_2]['data'] = {'month_data_2': event['month_data_2'][month_data_2]}
                else:
                    slide_idx[month_data_2]['data']['month_data_2'] = event['month_data_2'][month_data_2]
        
        # 3달전
        if 'month_data_3' in event:
            for month_data_3 in event['month_data_3']:
                agenda_list.append(month_data_3)
                if 'data' not in slide_idx[month_data_3]:
                    slide_idx[month_data_3]['data'] = {'month_data_3': event['month_data_3'][month_data_3]}
                else:
                    slide_idx[month_data_3]['data']['month_data_3'] = event['month_data_3'][month_data_3]
                    
        # 기타 데이터 추가
        slide_idx['agenda']['data'] = {'agenda_list': agenda_list}
        slide_idx['main']['data'] = {'customer_name': event['customer_name']}
        slide_idx['information']['data'] = {
            'customer_project': event['customer_project'],
            'account_id': event['account_id'],
            'manager_name': event['manager_name'],
            'msp_sales': event['msp_sales'],
            'msp_level': event['msp_level'],
            'report_cycle': event['report_cycle'],
            'customer_name': event['customer_name'],
            'region': event['region']
        }
        
        return slide_idx
    

class ConfigPpt:
    def presentation():
        prs = None
        s3 = boto3.client('s3')
        """ PPT 양식 파일 위치에 대한 정의"""
        form_object = s3.get_object(Bucket='FIX_INPUT: Bucket명', Key='FIX_INPUT: Report_Automation_4Month_form_V1_2.pptx')
        pptx_form = form_object['Body'].read()
        return Presentation(io.BytesIO(pptx_form))
    
    def index_ppt(prs):
        slide_idx = dict()
        for i, slide in enumerate(prs.slides):
            slide_idx[slide.slide_layout.name] = {
                'page_num': i,
                'region': 'no'
            }
            for k, shape in enumerate(slide.shapes):
                if shape.name == 'region':
                    slide_idx[slide.slide_layout.name]['region'] = 'yes'
                elif shape.has_chart:
                    if 'chart' not in slide_idx[slide.slide_layout.name]:
                        slide_idx[slide.slide_layout.name]['chart'] = {}
                    for t in shape.chart.series:
                        slide_idx[slide.slide_layout.name]['chart'][t.name] = t.values
                elif shape.has_table:
                    if 'table' not in slide_idx[slide.slide_layout.name]:
                        slide_idx[slide.slide_layout.name]['table'] = {}
                    slide_idx[slide.slide_layout.name]['table'][shape.name] = {
                        'max_x': len(shape.table.rows),
                        'max_y': len(shape.table.columns)
                    }
        return slide_idx
    
    def table_initialize(shape, start_x, start_y, max_x, max_y):
        for x in range(start_x, max_x):
            for y in range(start_y, max_y):
                shape.table.cell(x,y).text = ''
                ConfigPpt.table_font(shape.table.cell(x,y))
        return shape
    
    def rename_date(shape, month, x_axis, y_axis):
        shape.table.cell(x_axis ,y_axis).text = month[0]
        ConfigPpt.table_header_font(shape.table.cell(x_axis, y_axis))
        shape.table.cell(x_axis, y_axis+1).text = month[1]
        ConfigPpt.table_header_font(shape.table.cell(x_axis, y_axis+1))
        shape.table.cell(x_axis, y_axis+2).text = month[2]
        ConfigPpt.table_header_font(shape.table.cell(x_axis, y_axis+2))
        shape.table.cell(x_axis, y_axis+3).text = month[3]
        ConfigPpt.table_header_font(shape.table.cell(x_axis, y_axis+3))
        return shape
    
    def fetch_data(target, d1, d2=None, type='str', zerofill=None):
        if type == 'cmp':
            if d2 == None:
                d2 = d1
            result = d2 - d1
            if result > 0:
                target.text = str(format(int(d2), ',d'))+' (▲'+str(result)+')'
                ConfigPpt.text_red(target)
            elif result < 0:
                target.text = str(format(int(d2), ',d'))+' (▼'+str(result)+')'
                ConfigPpt.text_blue(target)
            elif result == 0:
                target.text = str(format(int(d2), ',d'))
                ConfigPpt.table_font(target)
                
        if type == 'onlycmp':
            if d2 == None:
                d2 = d1
            result = d2 - d1
            if result > 0:
                target.text = str(format(int(d2), ',d'))+' (▲)'
                ConfigPpt.text_red(target)
            elif result < 0:
                target.text = str(format(int(d2), ',d'))+' (▼)'
                ConfigPpt.text_blue(target)
            elif result == 0:
                target.text = str(format(int(d2), ',d'))
                ConfigPpt.table_font(target)
                
        elif type == 'str':
            if zerofill == 'no':
                if str(d1).isdigit():
                    if int(d1) == 0:
                        target.text = ''
                        ConfigPpt.table_font(target)
                        return
            target.text = str(d1)
            ConfigPpt.table_font(target)
    
    def table_header_font(target):
        target.text_frame.paragraphs[0].font.size = Pt(12)
        target.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        target.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        target.text_frame.paragraphs[0].font.bold = True
        return target
    
    def text_red_bold(target):
        ConfigPpt.text_bold(target)
        ConfigPpt.text_red(target)
        ConfigPpt.table_font(target)
    
    def text_red(target):
        target.text_frame.paragraphs[0].font.color.rgb = RGBColor(204, 0, 0)
        ConfigPpt.table_font(target)
    
    def text_blue(target):
        target.text_frame.paragraphs[0].font.color.rgb = RGBColor(47, 85, 151)
        ConfigPpt.table_font(target)
    
    def text_bold(target):
        target.text_frame.paragraphs[0].font.bold = True
    
    def table_font(target):
        target.text_frame.paragraphs[0].font.size = Pt(11)
        target.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    def copy_slide(slide_idx):
        # https://github.com/scanny/python-pptx/issues/68
        
        # 원본 레이아웃
        origin_slide = prs.slides[slide_idx]
        
        # 맨 끝부분에 레이아웃 복사 (PPT 슬라이드 번호 - 1)
        new_slide = prs.slides.add_slide(prs.slide_layouts[slide_idx])
        
        # Shape 복사
        for shape in origin_slide.shapes:
            # print(shape.name)
            newel = copy.deepcopy(shape.element)
            new_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')
        
        # Shape과 Relation
        for key, value in origin_slide.part.rels.items():
            if 'noteSlide' not in value.reltype:
                # if the relationship was a chart, we need to duplicate the embedded chart part and xlsx
                target = value._target
                if 'chart' in value.reltype:
                    partname = target.package.next_partname(
                        ChartPart.partname_template)
                    xlsx_blob = target.chart_workbook.xlsx_part.blob
                    target = ChartPart(
                        partname, target.content_type,
                        copy.deepcopy(target._element), package=target.package)
                    target.chart_workbook.xlsx_part = EmbeddedXlsxPart.new(
                        xlsx_blob, target.package)
                # Relation
                new_slide.part.rels.add_relationship(value.reltype, target,value.rId)
        
        # Slide 이동
        xml_slides = prs.slides._sldIdLst
        slides = list(xml_slides)
        xml_slides.remove(slides[-1])
        xml_slides.insert(slide_idx+1, slides[-1])
        return slide_idx + 1
    
    def hide_slide(slide_idx):
        del_slide = prs.slides[slide_idx]
        sld = del_slide._element
        sld.set('show', '0')
        return prs
    
    
class FixPage:
    def main(event, create_month, report_cycle):
        now_slide = prs.slides[event['page_num']]
        for s in now_slide.shapes:
            if s.name == 'customer_name':
                s.text = event['data']['customer_name']
                ConfigPpt.text_bold(s)
            elif s.name == 'report_cycle':
                s.text = report_cycle
            elif s.name == 'create_month':
                s.text = create_month
        return prs
    
    def report_intro_agenda(event):
        now_slide = prs.slides[event['page_num']]
        agenda_list = list(set(event['data']['agenda_list']))
        agenda_text = ['운영 정보']
        oper_list = [
            'other', 'elb', 
            'ec2', 'ebs',
            'rds_single', 'rds_cluster',
            'cache', 'cloud_front',
            'route53'
        ]
        
        if 'iam_summary' in agenda_list:
            agenda_text.append('IAM 현황')
        for title in agenda_list:
            if title in oper_list:
                agenda_text.append('자원 현황')
                break
        if 'work_summary' in agenda_list:
            agenda_text.append('업무 현황')
        if 'backup' in agenda_list:
            agenda_text.append('백업 현황')
        if 'ri_info' in agenda_list:
            agenda_text.append('요금 약정 현황')
        if 'trusted_advisor' in agenda_list:
            agenda_text.append('권고 사항')
        agenda_text.append('종합 의견')
        
        one_value = '\n'.join(agenda_text)
        
        for s in now_slide.shapes:
            if s.name == 'agenda':
                s.text = one_value
    
    def information(event, report_date):
        now_slide = prs.slides[event['page_num']]
        for s in now_slide.shapes:
            if s.name == 'customer_name':
                s.text = event['data']['customer_name']
            elif s.name == 'report_cycle':
                s.text = event['data']['report_cycle']
            elif s.name == 'msp_level':
                s.text = event['data']['msp_level']
            elif s.name == 'msp_sales':
                s.text = event['data']['msp_sales']
            elif s.name == 'manager_name':
                s.text = event['data']['manager_name']
            elif s.name == 'account_id':
                s.text = event['data']['account_id']
            elif s.name == 'customer_project':
                s.text = event['data']['customer_project']
            elif s.name == 'report_date':
                s.text = report_date
            elif s.name == 'region':
                s.text = ', '.join(event['data']['region'])
        return prs
    
    def iam_summary(event, month):
        if 'data' not in event:
            ConfigPpt.hide_slide(event['page_num'])
            return prs
        
        # 테이블 값 정리
        start_x = 1
        start_y = 1
        max_x = event['table']['iam_summary_table']['max_x']
        max_y = event['table']['iam_summary_table']['max_y']
        
        start_2nd_x = 1
        start_2nd_y = 0
        max_2nd_x = event['table']['iam_password_policy_table']['max_x']
        max_2nd_y = event['table']['iam_password_policy_table']['max_y']
        
        now_slide = prs.slides[event['page_num']]
        category = ['group', 'user', 'role', 'policy']
        
        # title을 key로 데이터 재생성
        formed_data = {}
        since_items = ['month_data_3', 'month_data_2', 'month_data_1', 'month_data_0']
        for since_item in since_items:
            if since_item in event['data']:
                for title in event['data'][since_item]:
                    if title not in formed_data:
                        formed_data[title] = {}
                    for data in event['data'][since_item][title]:
                        if data in formed_data[title]:
                            formed_data[title][data][since_item] = event['data'][since_item][title][data]
                        else:
                            formed_data[title][data] = {since_item: event['data'][since_item][title][data]}
                            
        for shape in now_slide.shapes:
            # 테이블
            if shape.name == 'iam_summary_table':
                ConfigPpt.table_initialize(shape, start_x, start_y, max_x, max_y)
                ConfigPpt.rename_date(shape, month, 0, 1)
                for x, item in enumerate(category, start=start_x):
                    ConfigPpt.fetch_data(shape.table.cell(x,1), formed_data['iam_summary'][item].get('month_data_3', 0), type='str')
                    ConfigPpt.fetch_data(shape.table.cell(x,2), formed_data['iam_summary'][item].get('month_data_3', 0), formed_data['iam_summary'][item].get('month_data_2', 0), type='cmp')
                    ConfigPpt.fetch_data(shape.table.cell(x,3), formed_data['iam_summary'][item].get('month_data_2', 0), formed_data['iam_summary'][item].get('month_data_1', 0), type='cmp')
                    ConfigPpt.fetch_data(shape.table.cell(x,4), formed_data['iam_summary'][item].get('month_data_1', 0), formed_data['iam_summary'][item].get('month_data_0', 0), type='cmp')
            
            elif shape.name == 'iam_password_policy_table':
                now_policy = event['data']['month_data_0']['iam_password_policy']
                shape.table.cell(2,1).text = now_policy['root_mfa_active']
                if len(now_policy) == 1:
                    for y in range(9):
                        y = y + 2
                        shape.table.cell(2,y).text = 'X'
                elif len(now_policy) > 1:
                    shape.table.cell(2,2).text = str(now_policy['minimum_password_length'])+'자'
                    shape.table.cell(2,3).text = now_policy['require_symbols']
                    shape.table.cell(2,4).text = now_policy['require_numbers']
                    shape.table.cell(2,5).text = now_policy['require_uppercase_characters']
                    shape.table.cell(2,6).text = now_policy['require_lowercase_characters']
                    shape.table.cell(2,7).text = now_policy['expire_passwords']
                    if now_policy['expire_passwords'] == 'X':
                        shape.table.cell(2,8).text = 'X'
                    else:
                        shape.table.cell(2,8).text = str(now_policy['max_password_age'])+'일'
                    shape.table.cell(2,9).text = now_policy['prevent_password_reuse_enable']
                    if now_policy['prevent_password_reuse_enable'] == 'X':
                        shape.table.cell(2,10).text = 'X'
                    else:
                        shape.table.cell(2,10).text = str(now_policy['password_reuse_prevention'])+'ea'
                for y in range(10):
                    y = y + 1
                    if 'X' in shape.table.cell(2,y).text:
                        ConfigPpt.text_red_bold(shape.table.cell(2,y))
                        ConfigPpt.table_font(shape.table.cell(2,y))
                    else:
                        ConfigPpt.text_bold(shape.table.cell(2,y))
                        ConfigPpt.table_font(shape.table.cell(2,y))
            
            # 차트
            elif shape.name == 'iam_summary_chart':
                
                chart_data = ChartData()
                chart_data.categories = ['Policy (정책)', 'Role (역할)', 'User (사용자)', 'Group (그룹)']
                
                chart_data.add_series(month[3], (
                    formed_data['iam_summary']['policy'].get('month_data_0', 0),
                    formed_data['iam_summary']['role'].get('month_data_0', 0),
                    formed_data['iam_summary']['user'].get('month_data_0', 0),
                    formed_data['iam_summary']['group'].get('month_data_0', 0)
                    )
                )
                chart_data.add_series(month[2], (
                    formed_data['iam_summary']['policy'].get('month_data_1', 0),
                    formed_data['iam_summary']['role'].get('month_data_1', 0),
                    formed_data['iam_summary']['user'].get('month_data_1', 0),
                    formed_data['iam_summary']['group'].get('month_data_1', 0)
                    )
                )
                chart_data.add_series(month[1], (
                    formed_data['iam_summary']['policy'].get('month_data_2', 0),
                    formed_data['iam_summary']['role'].get('month_data_2', 0),
                    formed_data['iam_summary']['user'].get('month_data_2', 0),
                    formed_data['iam_summary']['group'].get('month_data_2', 0)
                    )
                )
                chart_data.add_series(month[0], (
                    formed_data['iam_summary']['policy'].get('month_data_3', 0),
                    formed_data['iam_summary']['role'].get('month_data_3', 0),
                    formed_data['iam_summary']['user'].get('month_data_3', 0),
                    formed_data['iam_summary']['group'].get('month_data_3', 0)
                    )
                )
                shape.chart.replace_data(chart_data)
        
    def work_summary(event):
        if 'month_data_0' not in event['data']:
            ConfigPpt.hide_slide(event['page_num'])
            return prs
        else:
            this_event = event['data']['month_data_0']
            now_slide = prs.slides[event['page_num']]
            work_table = {}
            for x in range(4):
                x = x + 1
                for y in range(5):
                    y = y + 1
                    work_table[str(x)+'-'+str(y)] = 0
            for shape in now_slide.shapes:
                if shape.name == 'WK_chart' and len(this_event) > 0:
                    for work, value in this_event.items():
                        if 'Server' in work:
                            x = 1
                        elif 'Network' in work:
                            x = 2
                        elif 'Database' in work:
                            x = 3
                        elif 'ETC' in work:
                            x = 4
                        else:
                            x = 'None'
                        if '변경관리' in work:
                            y = 1
                        elif '보안관리' in work:
                            y = 2
                        elif '백업관리' in work:
                            y = 3
                        elif '문제관리' in work:
                            y = 4
                        elif '사고관리' in work:
                            y = 5
                        else:
                            y = 'None'
                        if x != 'None' or y != 'None':
                            work_table[str(x)+'-'+str(y)] = value
                    work_chart = ChartData()
                    work_chart.categories = ['변경관리', '보안관리', '백업관리', '문제관리', '사고관리']
                    work_chart.add_series('Server', (work_table['1-1'],
                        work_table['1-2'], work_table['1-3'],
                        work_table['1-4'], work_table['1-5']))
                    work_chart.add_series('Network', (work_table['2-1'],
                        work_table['2-2'], work_table['2-3'],
                        work_table['2-4'], work_table['2-5']))
                    work_chart.add_series('Database', (work_table['3-1'],
                        work_table['3-2'], work_table['3-3'],
                        work_table['3-4'], work_table['3-5']))
                    work_chart.add_series('ETC', (work_table['4-1'],
                        work_table['4-2'], work_table['4-3'],
                        work_table['4-4'], work_table['4-5']))
                    shape.chart.replace_data(work_chart)
            return prs
    
    
class ElasticPage:
    def trusted_advisor(event):
        # 페이지 위치 / 데이터가 없을 경우 숨기기
        page_num = event['page_num']
        if 'data' not in event:
            ConfigPpt.hide_slide(page_num)
            return prs
        else:
            if 'month_data_0' not in event['data']:
                ConfigPpt.hide_slide(page_num)
                return prs
                
        # 테이블 값 정리
        start_x = 1
        start_y = 0
        max_x = event['table']['TA_table']['max_x'] # 6
        max_y = event['table']['TA_table']['max_y'] # 3
        repeat_max_counter = 0
        repeat_counter = 0
        
        # 2차원 배열로 변경
        formed_data = {}
        for ta_category in event['data']['month_data_0']:
            if ta_category not in formed_data:
                formed_data[ta_category] = []
            for num in range(math.ceil(len(event['data']['month_data_0'][ta_category])/(max_x-start_x))):
                formed_data[ta_category].append(event['data']['month_data_0'][ta_category][num*(max_x-start_x) : (max_x-start_x)+num*(max_x-start_x)])
                repeat_max_counter += 1
        
        # 데이터 삽입
        category_list = {'security': '보안', 'fault_tolerance': '내결함성', 'performance': '성능', 'service_limits': '한도', 'cost_optimization': '비용 최적화'}
        for ta_category in formed_data:
            for data_list in formed_data[ta_category]:
                now_slide = prs.slides[page_num]
                repeat_counter += 1
                for shape in now_slide.shapes:
                    if shape.name == 'Text Placeholder 1':
                        now_slide.shapes.element.remove(shape.element)
                        
                    elif shape.name == 'TA_category':
                        shape.text = category_list.get(ta_category)
                    
                    elif shape.name == 'TA_table':
                        ConfigPpt.table_initialize(shape, start_x, start_y, max_x, max_y)
                        for x, item in enumerate(data_list, start=start_x):
                            
                            ConfigPpt.fetch_data(shape.table.cell(x, 0), item['name'])
                            ConfigPpt.fetch_data(shape.table.cell(x, 1), item['description'])
                            formed_result = ", ".join(item['result'][0:3]) + ' 외 ' + str(len(item['result'][2:])) + 'ea' if len(item['result']) > 3 else ", ".join(item['result'])
                            ConfigPpt.fetch_data(shape.table.cell(x, 2), formed_result)
                            
                            # 글씨 속성 변경
                            ConfigPpt.text_bold(shape.table.cell(x,0))
                    
                # 페이지 복사
                if repeat_max_counter > repeat_counter:
                    page_num = ConfigPpt.copy_slide(page_num)    

    def ri_usage(event):
        # 페이지 위치 / 데이터가 없을 경우 숨기기
        page_num = event['page_num']
        if 'data' not in event:
            ConfigPpt.hide_slide(page_num)
            return prs
        else:
            if 'month_data_0' not in event['data']:
                ConfigPpt.hide_slide(page_num)
                return prs
        
        # 테이블 값 정리
        start_x = 1
        start_y = 0
        max_x = event['table']['RI_used_table']['max_x'] # 6
        max_y = event['table']['RI_used_table']['max_y'] # 3
        repeat_max_counter = 0
        repeat_counter = 0
        
        # 2차원 배열로 변경
        formed_data = {}
        for region in event['data']['month_data_0']:
            if region not in formed_data:
                formed_data[region] = []
            for num in range(math.ceil(len(event['data']['month_data_0'][region])/(max_x-start_x))):
                formed_data[region].append(event['data']['month_data_0'][region][num*(max_x-start_x) : (max_x-start_x)+num*(max_x-start_x)])
                repeat_max_counter += 1
            
        # 데이터 삽입
        for region in formed_data:
            for data_list in formed_data[region]:
                now_slide = prs.slides[page_num]
                repeat_counter += 1
                for shape in now_slide.shapes:
                    if shape.name == 'Text Placeholder 1':
                        now_slide.shapes.element.remove(shape.element)
                        
                    elif shape.name == 'region':
                        shape.text = region
                        
                    elif shape.name == 'RI_used_table':
                        ConfigPpt.table_initialize(shape, start_x, start_y, max_x, max_y)
                        for x, item in enumerate(data_list, start=start_x):
                            ConfigPpt.fetch_data(shape.table.cell(x, 0), item['service'])
                            ConfigPpt.fetch_data(shape.table.cell(x, 1), item['platform'])
                            ConfigPpt.fetch_data(shape.table.cell(x, 2), item['instance_type'])
                            ConfigPpt.fetch_data(shape.table.cell(x, 3), item['normalized_used_percent'])
                            
                            # 글씨 속성 변경
                            ConfigPpt.text_bold(shape.table.cell(x,0))
                            if int(item['normalized_used_percent'].split('%')[0]) < 60:
                                ConfigPpt.text_red_bold(shape.table.cell(x,3))
                
                # 페이지 복사
                if repeat_max_counter > repeat_counter:
                    page_num = ConfigPpt.copy_slide(page_num)
    
    def ri_info(event):
        # 페이지 위치 / 데이터가 없을 경우 숨기기
        page_num = event['page_num']
        if 'data' not in event:
            ConfigPpt.hide_slide(page_num)
            return prs
        else:
            if 'month_data_0' not in event['data']:
                ConfigPpt.hide_slide(page_num)
                return prs
        
        # 테이블 값 정리
        start_x = 1
        start_y = 0
        max_x = event['table']['RI_table']['max_x']
        max_y = event['table']['RI_table']['max_y']
        repeat_max_counter = 0
        repeat_counter = 0
        
        # 2차원 배열로 변경
        formed_data = {}
        for region in event['data']['month_data_0']:
            if region not in formed_data:
                formed_data[region] = []
            for num in range(math.ceil(len(event['data']['month_data_0'][region])/(max_x-start_x))):
                formed_data[region].append(event['data']['month_data_0'][region][num*(max_x-start_x) : (max_x-start_x)+num*(max_x-start_x)])
                repeat_max_counter += 1
        
        # 데이터 삽입
        for region in formed_data:
            for data_list in formed_data[region]:
                now_slide = prs.slides[page_num]
                repeat_counter += 1
                for shape in now_slide.shapes:
                    if shape.name == 'Text Placeholder 1':
                        now_slide.shapes.element.remove(shape.element)
                        
                    elif shape.name == 'region':
                        shape.text = region
                        
                    elif shape.name == 'RI_table':
                        ConfigPpt.table_initialize(shape, start_x, start_y, max_x, max_y)
                        for x, item in enumerate(data_list, start=start_x):
                            ConfigPpt.fetch_data(shape.table.cell(x, 0), item['service'])
                            ConfigPpt.fetch_data(shape.table.cell(x, 1), item['platform'])
                            ConfigPpt.fetch_data(shape.table.cell(x, 2), item['instance_type'])
                            ConfigPpt.fetch_data(shape.table.cell(x, 3), item['count'])
                            ConfigPpt.fetch_data(shape.table.cell(x, 4), item['offering_class'])
                            ConfigPpt.fetch_data(shape.table.cell(x, 5), item['offering_type'])
                            ConfigPpt.fetch_data(shape.table.cell(x, 6), item['start'])
                            ConfigPpt.fetch_data(shape.table.cell(x, 7), item['end'])
                            
                            # 글씨 속성 변경
                            ConfigPpt.text_bold(shape.table.cell(x,0))
                
                # 페이지 복사
                if repeat_max_counter > repeat_counter:
                    page_num = ConfigPpt.copy_slide(page_num)
    
    def backup(event):
        # 페이지 위치 / 데이터가 없을 경우 숨기기
        page_num = event['page_num']
        if 'data' not in event:
            ConfigPpt.hide_slide(page_num)
            return prs
        else:
            if 'month_data_0' not in event['data']:
                ConfigPpt.hide_slide(page_num)
                return prs
        
        # 테이블 값 정리
        start_x = 1
        start_y = 0
        max_x = event['table']['BU_table']['max_x']
        max_y = event['table']['BU_table']['max_y']
        repeat_max_counter = 0
        repeat_counter = 0
        
        # 2차원 배열로 변경
        formed_data = {}
        for region in event['data']['month_data_0']:
            if region not in formed_data:
                formed_data[region] = []
            for num in range(math.ceil(len(event['data']['month_data_0'][region])/(max_x-start_x))):
                formed_data[region].append(event['data']['month_data_0'][region][num*(max_x-start_x) : (max_x-start_x)+num*(max_x-start_x)])
                repeat_max_counter += 1
        
        # 데이터 삽입
        for region in formed_data:
            for data_list in formed_data[region]:
                now_slide = prs.slides[page_num]
                repeat_counter += 1
                for shape in now_slide.shapes:
                    if shape.name == 'Text Placeholder 1':
                        now_slide.shapes.element.remove(shape.element)
                        
                    elif shape.name == 'region':
                        shape.text = region
                        
                    elif shape.name == 'BU_table':
                        ConfigPpt.table_initialize(shape, start_x, start_y, max_x, max_y)
                        for x, item in enumerate(data_list, start=start_x):
                            ConfigPpt.fetch_data(shape.table.cell(x, 0), item['source'])
                            ConfigPpt.fetch_data(shape.table.cell(x, 1), item['service'])
                            ConfigPpt.fetch_data(shape.table.cell(x, 2), item['converted_schedule_expression'])
                            ConfigPpt.fetch_data(shape.table.cell(x, 3), item['life_cycle'])
                            ConfigPpt.fetch_data(shape.table.cell(x, 4), item['target'])
                            ConfigPpt.fetch_data(shape.table.cell(x, 5), item['count'])
                            
                            # 글씨 속성 변경
                            ConfigPpt.text_bold(shape.table.cell(x,0))
                            
                # 페이지 복사
                if repeat_max_counter > repeat_counter:
                    page_num = ConfigPpt.copy_slide(page_num)
    
    def work_list(event):
        # 페이지 위치 / 데이터가 없을 경우 숨기기
        page_num = event['page_num']
        if 'data' not in event:
            ConfigPpt.hide_slide(page_num)
            return prs
        else:
            if 'month_data_0' not in event['data']:
                ConfigPpt.hide_slide(page_num)
                return prs
        
        # 테이블 값 정리
        start_x = 1
        start_y = 0
        max_x = event['table']['WK_table']['max_x']
        max_y = event['table']['WK_table']['max_y']
        
        # 2차원 배열로 변경
        formed_data = []
        for num in range(math.ceil(len(event['data']['month_data_0'])/(max_x-start_x))):
            formed_data.append(event['data']['month_data_0'][num*(max_x-start_x) : (max_x-start_x)+num*(max_x-start_x)])
        
        # 데이터 삽입
        for list_count, data_list in enumerate(formed_data):
            now_slide = prs.slides[page_num]
            for shape in now_slide.shapes:
                if shape.name == 'Text Placeholder 1':
                    now_slide.shapes.element.remove(shape.element)
                    
                elif shape.name == 'WK_table':
                    ConfigPpt.table_initialize(shape, start_x, start_y, max_x, max_y)
                    for x, item in enumerate(data_list, start=start_x):
                        ConfigPpt.fetch_data(shape.table.cell(x, 0), item['work_type'])
                        ConfigPpt.fetch_data(shape.table.cell(x, 1), item['service'])
                        ConfigPpt.fetch_data(shape.table.cell(x, 2), item['work_date'])
                        ConfigPpt.fetch_data(shape.table.cell(x, 3), item['work_value'])
                        
                        # 글씨 속성 변경
                        ConfigPpt.text_bold(shape.table.cell(x,0))
            
            # 페이지 복사
            if len(formed_data) > list_count+1:
                page_num = ConfigPpt.copy_slide(page_num)
    
    def other(event, month):
        # 페이지 위치 / 데이터가 없을 경우 숨기기
        page_num = event['page_num']
        if 'data' not in event:
            ConfigPpt.hide_slide(page_num)
            return prs
        
        # 테이블 값 정리
        start_x = 1
        start_y = 0
        max_x = event['table']['Other_table']['max_x']
        max_y = event['table']['Other_table']['max_y']
        
        # $since_items들 뭉치기
        temp_data = {}
        for since_item in event['data']:
            for item in event['data'][since_item]:
                if item['region'] + '_' + item['service'] + '_' + item['category'] not in temp_data:
                    temp_data[item['region'] + '_' + item['service'] + '_' + item['category']] \
                    = {'region': item['region'], 'service': item['service'], 'category': item['category']}
                temp_data[item['region'] + '_' + item['service'] + '_' + item['category']][since_item] = item['count']
        
        temp_data = list(map(lambda x:temp_data[x], temp_data)) 
        
        # 2차원 배열로 변경 / 더이상 사용하지 않는 temp_data 제거
        formed_data = []
        for num in range(math.ceil(len(temp_data)/(max_x-start_x))):
            formed_data.append(temp_data[num*(max_x-start_x) : (max_x-start_x)+num*(max_x-start_x)])
        del temp_data
        
        # 반복횟수 = 2차원 배열 행 수 
        for list_count, data_list in enumerate(formed_data):
            now_slide = prs.slides[page_num]
            for shape in now_slide.shapes:
                if shape.name == 'Text Placeholder 1':
                    now_slide.shapes.element.remove(shape.element)
                    
                elif shape.name == 'Other_table':
                    ConfigPpt.table_initialize(shape, start_x, start_y, max_x, max_y)
                    ConfigPpt.rename_date(shape, month, 0, 3)
                    for x, item in enumerate(data_list, start=start_x):
                        ConfigPpt.fetch_data(shape.table.cell(x, 0), item['region'])
                        ConfigPpt.fetch_data(shape.table.cell(x, 1), item['service'])
                        ConfigPpt.fetch_data(shape.table.cell(x, 2), item['category'])
                        ConfigPpt.fetch_data(shape.table.cell(x, 3), item.get('month_data_3', 0))
                        ConfigPpt.fetch_data(shape.table.cell(x, 4), item.get('month_data_3', 0), item.get('month_data_2', 0), type='cmp')
                        ConfigPpt.fetch_data(shape.table.cell(x, 5), item.get('month_data_2', 0), item.get('month_data_1', 0), type='cmp')
                        ConfigPpt.fetch_data(shape.table.cell(x, 6), item.get('month_data_1', 0), item.get('month_data_0', 0), type='cmp')
                        
            # 페이지 복사
            if len(formed_data) > list_count+1:
                page_num = ConfigPpt.copy_slide(page_num)
    
    def route53(event):
        # 페이지 위치 / 데이터가 없을 경우 숨기기
        page_num = event['page_num']
        if 'data' not in event:
            ConfigPpt.hide_slide(page_num)
            return prs
        else:
            if 'month_data_0' not in event['data']:
                ConfigPpt.hide_slide(page_num)
                return prs
        
        # 테이블 값 정리
        start_x = 2
        start_y = 0
        max_x = event['table']['R53_table']['max_x']
        max_y = event['table']['R53_table']['max_y']
        
        # 사용하지 않는 key 제거
        event['data'].pop('month_data_3', None)
        event['data'].pop('month_data_2', None)
        event['data'].pop('month_data_1', None)
        
        # 2차원 배열로 변경
        formed_data = []
        for num in range(math.ceil(len(event['data']['month_data_0'])/(max_x-start_x))):
            formed_data.append(event['data']['month_data_0'][num*(max_x-start_x) : (max_x-start_x)+num*(max_x-start_x)])
        
        # 반복횟수 = 2차원 배열 행 수 
        for list_count, data_list in enumerate(formed_data):
            now_slide = prs.slides[page_num]
            for shape in now_slide.shapes:
                if shape.name == 'Text Placeholder 1':
                    now_slide.shapes.element.remove(shape.element)
                elif shape.name == 'R53_table':
                    ConfigPpt.table_initialize(shape, start_x, start_y, max_x, max_y)
                    for x, item in enumerate(data_list, start=start_x):
                        ConfigPpt.fetch_data(shape.table.cell(x, 0), (max_x-start_x)*list_count + (x-1))
                        ConfigPpt.fetch_data(shape.table.cell(x, 1), item['domain_name'])
                        ConfigPpt.fetch_data(shape.table.cell(x, 2), item['type'])
                        ConfigPpt.fetch_data(shape.table.cell(x, 3), item.get('total', '0'))
                        ConfigPpt.fetch_data(shape.table.cell(x, 4), item.get('record').get('TXT', '0'))
                        ConfigPpt.fetch_data(shape.table.cell(x, 5), item.get('record').get('CNAME', '0'))
                        ConfigPpt.fetch_data(shape.table.cell(x, 6), item.get('record').get('A', '0'))
                        ConfigPpt.fetch_data(shape.table.cell(x, 7), item.get('record').get('NS', '0'))
                        ConfigPpt.fetch_data(shape.table.cell(x, 8), item.get('record').get('ETC', '0'))
                        
                        # 글씨 속성 변경
                        for y in range(2):
                            ConfigPpt.text_bold(shape.table.cell(x,y))
                            
            # 페이지 복사
            if len(formed_data) > list_count+1:
                page_num = ConfigPpt.copy_slide(page_num)
    
    def cloudfront(event, month):
        # 페이지 위치 / 데이터가 없을 경우 숨기기
        page_num = event['page_num']
        if 'data' not in event:
            ConfigPpt.hide_slide(page_num)
            return prs
        
        # 테이블 값 정리
        start_x = 2
        start_y = 0
        max_x = event['table']['CF_table']['max_x']
        max_y = event['table']['CF_table']['max_y']
        
        # $since_items들 뭉치기
        temp_data = {}
        for since_item in event['data']:
            for item in event['data'][since_item]:
                if item['domain_name'] not in temp_data:
                    temp_data[item['domain_name']] = {'domain_name': item['domain_name'], 'cname': item['cname'], 'price_class': item['price_class'], 'status': item['status'], \
                    'month_data_3': 0, 'month_data_2': 0, 'month_data_1': 0, 'month_data_0': 0}
                # temp_data[item['domain_name']][since_item] = str(format(item['request'], ',d'))
                temp_data[item['domain_name']][since_item] = item['request']
        
        # 2차원 배열로 변경하기 위해 dict => list로 변경
        temp_data = list(map(lambda x:temp_data[x], temp_data))
        
        # 2차원 배열로 변경 / 사용하지 않는 temp_data 제거
        formed_data = []
        for num in range(math.ceil(len(temp_data)/(max_x-start_x))):
            formed_data.append(temp_data[num*(max_x-start_x) : (max_x-start_x)+num*(max_x-start_x)])
        del temp_data
        
        # 반복횟수 = 2차원 배열 행 수 
        for list_count, data_list in enumerate(formed_data):
            now_slide = prs.slides[page_num]
            for shape in now_slide.shapes:
                if shape.name == 'Text Placeholder 1':
                    now_slide.shapes.element.remove(shape.element)
                elif shape.name == 'CF_table':
                    ConfigPpt.table_initialize(shape, start_x, start_y, max_x, max_y)
                    ConfigPpt.rename_date(shape, month, 1, 3)
                    for x, item in enumerate(data_list, start=start_x):
                        ConfigPpt.fetch_data(shape.table.cell(x, 0), item['domain_name'])
                        ConfigPpt.fetch_data(shape.table.cell(x, 1), item['price_class'])
                        ConfigPpt.fetch_data(shape.table.cell(x, 2), item['cname'])
                        ConfigPpt.fetch_data(shape.table.cell(x, 3), item['month_data_3'], type='onlycmp')
                        ConfigPpt.fetch_data(shape.table.cell(x, 4), item['month_data_3'], item['month_data_2'], type='onlycmp')
                        ConfigPpt.fetch_data(shape.table.cell(x, 5), item['month_data_2'], item['month_data_1'], type='onlycmp')
                        ConfigPpt.fetch_data(shape.table.cell(x, 6), item['month_data_1'], item['month_data_0'], type='onlycmp')
                        
                        # 글씨 속성 변경
                        for y in range(1):
                            ConfigPpt.text_bold(shape.table.cell(x,y))
            # 페이지 복사
            if len(formed_data) > list_count+1:
                page_num = ConfigPpt.copy_slide(page_num)
    
    def cache(event, month):
        # 페이지 위치 / 데이터가 없을 경우 숨기기
        page_num = event['page_num']
        if 'data' not in event:
            ConfigPpt.hide_slide(page_num)
            return prs
        
        # 테이블 값 정리
        start_x = 2
        start_y = 0
        max_x = event['table']['Cache_table']['max_x']
        max_y = event['table']['Cache_table']['max_y']
        
        start_2nd_x = 1
        start_2nd_y = 0
        max_2nd_x = event['table']['Cache_type_table']['max_x']
        max_2nd_y = event['table']['Cache_type_table']['max_y']
        
        # region을 key로 데이터 재생성
        formed_data = ConfigData.event_format(event, region=True)
        
        formed_data = ConfigData.data_maxlen_format(formed_data, region=True, key='cache_count', maxlen=max_x-start_x)
        formed_data = ConfigData.data_maxlen_format(formed_data, region=True, key='cache_type', maxlen=max_2nd_x-start_2nd_x)
        
        # 슬라이드 개수 = region개수 
        for region_count, region in enumerate(list(formed_data), start=1):
            now_slide = prs.slides[page_num]
            # 슬라이드의 각 shape수정
            for shape in now_slide.shapes:
                if shape.name == 'Text Placeholder 1':
                    now_slide.shapes.element.remove(shape.element)
                    
                elif shape.name == 'region':
                    shape.text = region
                    
                elif shape.name == 'Cache_chart':
                    chart_data = ChartData()
                    chart_data.categories = list(formed_data[region]['cache_count'])
                    chart_data.add_series(month[0], tuple(x['month_data_3'] for x in formed_data[region]['cache_count'].values()))
                    chart_data.add_series(month[1], tuple(x['month_data_2'] for x in formed_data[region]['cache_count'].values()))
                    chart_data.add_series(month[2], tuple(x['month_data_1'] for x in formed_data[region]['cache_count'].values()))
                    chart_data.add_series(month[3], tuple(x['month_data_0'] for x in formed_data[region]['cache_count'].values()))
                    shape.chart.replace_data(chart_data)
                    
                elif shape.name == 'Cache_type_chart':
                    chart_data = ChartData()
                    chart_data.categories = list(formed_data[region]['cache_type'])
                    chart_data.add_series(month[3], tuple(x['month_data_0'] for x in formed_data[region]['cache_type'].values()))
                    shape.chart.replace_data(chart_data)
                    
                elif shape.name == 'Cache_table':
                    # Initialize / 날짜 삽입
                    ConfigPpt.table_initialize(shape, start_x, start_y, max_x, max_y)
                    ConfigPpt.rename_date(shape, month, 1, 1)
                    
                    # 데이터 삽입
                    for x, item in enumerate(list(formed_data[region]['cache_count']), start=start_x):
                        ConfigPpt.fetch_data(shape.table.cell(x, 0), item)
                        ConfigPpt.fetch_data(shape.table.cell(x, 1), formed_data[region]['cache_count'][item]['month_data_3'])
                        ConfigPpt.fetch_data(shape.table.cell(x, 2), formed_data[region]['cache_count'][item]['month_data_3'], formed_data[region]['cache_count'][item]['month_data_2'], type='cmp')
                        ConfigPpt.fetch_data(shape.table.cell(x, 3), formed_data[region]['cache_count'][item]['month_data_2'], formed_data[region]['cache_count'][item]['month_data_1'], type='cmp')
                        ConfigPpt.fetch_data(shape.table.cell(x, 4), formed_data[region]['cache_count'][item]['month_data_1'], formed_data[region]['cache_count'][item]['month_data_0'], type='cmp')
                                
                elif shape.name == 'Cache_type_table':
                    # Initialize / 날짜 삽입
                    ConfigPpt.table_initialize(shape, start_2nd_x, start_2nd_y, max_2nd_x, max_2nd_y)
                    ConfigPpt.rename_date(shape, month, 0, 1)
                    
                    # 데이터 삽입
                    for x, item in enumerate(list(formed_data[region]['cache_type']), start=start_2nd_x):
                        ConfigPpt.fetch_data(shape.table.cell(x, 0), item)
                        ConfigPpt.fetch_data(shape.table.cell(x, 1), formed_data[region]['cache_type'][item]['month_data_3'])
                        ConfigPpt.fetch_data(shape.table.cell(x, 2), formed_data[region]['cache_type'][item]['month_data_3'], formed_data[region]['cache_type'][item]['month_data_2'], type='cmp')
                        ConfigPpt.fetch_data(shape.table.cell(x, 3), formed_data[region]['cache_type'][item]['month_data_2'], formed_data[region]['cache_type'][item]['month_data_1'], type='cmp')
                        ConfigPpt.fetch_data(shape.table.cell(x, 4), formed_data[region]['cache_type'][item]['month_data_1'], formed_data[region]['cache_type'][item]['month_data_0'], type='cmp')
                        
            # 페이지 복사
            if len(formed_data) > region_count:
                page_num = ConfigPpt.copy_slide(page_num)
    
    def rds_single(event, month):
        # 페이지 위치 / 데이터가 없을 경우 숨기기
        page_num = event['page_num']
        if 'data' not in event:
            ConfigPpt.hide_slide(page_num)
            return prs
        
        # 테이블 값 정리
        start_x = 1
        start_y = 0
        max_x = event['table']['RDS_table']['max_x']
        max_y = event['table']['RDS_table']['max_y']
        
        start_2nd_x = 1
        start_2nd_y = 0
        max_2nd_x = event['table']['RDS_type_table']['max_x']
        max_2nd_y = event['table']['RDS_type_table']['max_y']
        
        # region을 key로 데이터 재생성
        formed_data = ConfigData.event_format(event, region=True)
        formed_data = ConfigData.data_maxlen_format(formed_data, region=True, key='engines', maxlen=max_x-start_x)
        formed_data = ConfigData.data_maxlen_format(formed_data, region=True, key='db_class', maxlen=max_2nd_x-start_2nd_x)
        
        # 슬라이드 개수 = region개수 
        for region_count, region in enumerate(list(formed_data), start=1):
            now_slide = prs.slides[page_num]
            # 슬라이드의 각 shape수정
            for shape in now_slide.shapes:
                if shape.name == 'Text Placeholder 1':
                    now_slide.shapes.element.remove(shape.element)
                    
                elif shape.name == 'region':
                    shape.text = region
                    
                elif shape.name == 'RDS_chart':
                    chart_data = ChartData()
                    chart_data.categories = list(formed_data[region]['engines'])
                    chart_data.add_series(month[0], tuple(x['month_data_3'] for x in formed_data[region]['engines'].values()))
                    chart_data.add_series(month[1], tuple(x['month_data_2'] for x in formed_data[region]['engines'].values()))
                    chart_data.add_series(month[2], tuple(x['month_data_1'] for x in formed_data[region]['engines'].values()))
                    chart_data.add_series(month[3], tuple(x['month_data_0'] for x in formed_data[region]['engines'].values()))
                    shape.chart.replace_data(chart_data)
                    
                elif shape.name == 'RDS_type_chart':
                    chart_data = ChartData()
                    chart_data.categories = list(formed_data[region]['db_class'])
                    chart_data.add_series(month[3], tuple(x['month_data_0'] for x in formed_data[region]['db_class'].values()))
                    shape.chart.replace_data(chart_data)
                    
                elif shape.name == 'RDS_table':
                    # Initialize / 날짜 삽입
                    ConfigPpt.table_initialize(shape, start_x, start_y, max_x, max_y)
                    ConfigPpt.rename_date(shape, month, 0, 1)
                    
                    # 데이터 삽입
                    for x, item in enumerate(list(formed_data[region]['engines']), start=start_x):
                        ConfigPpt.fetch_data(shape.table.cell(x, 0), item)
                        ConfigPpt.fetch_data(shape.table.cell(x, 1), formed_data[region]['engines'][item]['month_data_3'])
                        ConfigPpt.fetch_data(shape.table.cell(x, 2), formed_data[region]['engines'][item]['month_data_3'], formed_data[region]['engines'][item]['month_data_2'], type='cmp')
                        ConfigPpt.fetch_data(shape.table.cell(x, 3), formed_data[region]['engines'][item]['month_data_2'], formed_data[region]['engines'][item]['month_data_1'], type='cmp')
                        ConfigPpt.fetch_data(shape.table.cell(x, 4), formed_data[region]['engines'][item]['month_data_1'], formed_data[region]['engines'][item]['month_data_0'], type='cmp')
                                
                elif shape.name == 'RDS_type_table':
                    # Initialize / 날짜 삽입
                    ConfigPpt.table_initialize(shape, start_2nd_x, start_2nd_y, max_2nd_x, max_2nd_y)
                    ConfigPpt.rename_date(shape, month, 0, 1)
                    
                    # 데이터 삽입
                    for x, item in enumerate(list(formed_data[region]['db_class']), start=start_2nd_x):
                        ConfigPpt.fetch_data(shape.table.cell(x, 0), item)
                        ConfigPpt.fetch_data(shape.table.cell(x, 1), formed_data[region]['db_class'][item]['month_data_3'])
                        ConfigPpt.fetch_data(shape.table.cell(x, 2), formed_data[region]['db_class'][item]['month_data_3'], formed_data[region]['db_class'][item]['month_data_2'], type='cmp')
                        ConfigPpt.fetch_data(shape.table.cell(x, 3), formed_data[region]['db_class'][item]['month_data_2'], formed_data[region]['db_class'][item]['month_data_1'], type='cmp')
                        ConfigPpt.fetch_data(shape.table.cell(x, 4), formed_data[region]['db_class'][item]['month_data_1'], formed_data[region]['db_class'][item]['month_data_0'], type='cmp')
                        
            # 페이지 복사
            if len(formed_data) > region_count:
                page_num = ConfigPpt.copy_slide(page_num)
    
    def rds_cluster(event, month):
        # 페이지 위치 / 데이터가 없을 경우 숨기기
        page_num = event['page_num']
        if 'data' not in event:
            ConfigPpt.hide_slide(page_num)
            return prs
        
        # 테이블 값 정리
        start_x = 1
        start_y = 0
        max_x = event['table']['ARR_table']['max_x']
        max_y = event['table']['ARR_table']['max_y']
        
        start_2nd_x = 1
        start_2nd_y = 0
        max_2nd_x = event['table']['ARR_type_table']['max_x']
        max_2nd_y = event['table']['ARR_type_table']['max_y']
        
        # $engine_$role: int 형태로 변형
        for since_item in event['data']:
            for region in event['data'][since_item]:
                temp_data = {}
                if 'engines' in event['data'][since_item][region]:
                    for engine in event['data'][since_item][region]['engines']:
                        for role in event['data'][since_item][region]['engines'][engine]:
                            temp_data[engine + '_' + role] = event['data'][since_item][region]['engines'][engine][role]
                event['data'][since_item][region].pop('engines')
                event['data'][since_item][region]['engines'] = temp_data
        
        # region을 key로 데이터 재생성                        
        formed_data = ConfigData.event_format(event, region=True)
        formed_data = ConfigData.data_maxlen_format(formed_data, region=True, key='engines', maxlen=max_x-start_x)
        formed_data = ConfigData.data_maxlen_format(formed_data, region=True, key='db_class', maxlen=max_2nd_x-start_2nd_x)
        
        # 슬라이드 개수 = region개수 
        for region_count, region in enumerate(list(formed_data), start=1):
            now_slide = prs.slides[page_num]
            # 슬라이드의 각 shape수정
            for shape in now_slide.shapes:
                if shape.name == 'Text Placeholder 1':
                    now_slide.shapes.element.remove(shape.element)
                    
                elif shape.name == 'region':
                    shape.text = region
                    
                elif shape.name == 'ARR_chart':
                    chart_data = ChartData()
                    chart_data.categories = list(formed_data[region]['engines'])
                    chart_data.add_series(month[0], tuple(x['month_data_3'] for x in formed_data[region]['engines'].values()))
                    chart_data.add_series(month[1], tuple(x['month_data_2'] for x in formed_data[region]['engines'].values()))
                    chart_data.add_series(month[2], tuple(x['month_data_1'] for x in formed_data[region]['engines'].values()))
                    chart_data.add_series(month[3], tuple(x['month_data_0'] for x in formed_data[region]['engines'].values()))
                    shape.chart.replace_data(chart_data)
                    
                elif shape.name == 'ARR_type_chart':
                    # serverless만 있을 때 타입 없음
                    if len(formed_data[region]['db_class']) == 0:
                        continue
                    chart_data = ChartData()
                    chart_data.categories = list(formed_data[region]['db_class'])
                    chart_data.add_series(month[3], tuple(x['month_data_0'] for x in formed_data[region]['db_class'].values()))
                    shape.chart.replace_data(chart_data)
                    
                elif shape.name == 'ARR_table':
                    # Initialize / 날짜 삽입
                    ConfigPpt.table_initialize(shape, start_x, start_y, max_x, max_y)
                    ConfigPpt.rename_date(shape, month, 0, 2)
                    
                    # 데이터 삽입
                    for x, item in enumerate(list(formed_data[region]['engines']), start=start_x):
                        # document DB 테스트 해야함
                        ConfigPpt.fetch_data(shape.table.cell(x, 0), item.split('_')[0])
                        ConfigPpt.fetch_data(shape.table.cell(x, 1), item.split('_')[1])
                        ConfigPpt.fetch_data(shape.table.cell(x, 2), formed_data[region]['engines'][item]['month_data_3'])
                        ConfigPpt.fetch_data(shape.table.cell(x, 3), formed_data[region]['engines'][item]['month_data_3'], formed_data[region]['engines'][item]['month_data_2'], type='cmp')
                        ConfigPpt.fetch_data(shape.table.cell(x, 4), formed_data[region]['engines'][item]['month_data_2'], formed_data[region]['engines'][item]['month_data_1'], type='cmp')
                        ConfigPpt.fetch_data(shape.table.cell(x, 5), formed_data[region]['engines'][item]['month_data_1'], formed_data[region]['engines'][item]['month_data_0'], type='cmp')
                        
                elif shape.name == 'ARR_type_table':
                    # Initialize / 날짜 삽입
                    ConfigPpt.table_initialize(shape, start_2nd_x, start_2nd_y, max_2nd_x, max_2nd_y)
                    ConfigPpt.rename_date(shape, month, 0, 1)
                    
                    # 데이터 삽입
                    for x, item in enumerate(list(formed_data[region]['db_class']), start=start_2nd_x):
                        ConfigPpt.fetch_data(shape.table.cell(x, 0), item)
                        ConfigPpt.fetch_data(shape.table.cell(x, 1), formed_data[region]['db_class'][item]['month_data_3'])
                        ConfigPpt.fetch_data(shape.table.cell(x, 2), formed_data[region]['db_class'][item]['month_data_3'], formed_data[region]['db_class'][item]['month_data_2'], type='cmp')
                        ConfigPpt.fetch_data(shape.table.cell(x, 3), formed_data[region]['db_class'][item]['month_data_2'], formed_data[region]['db_class'][item]['month_data_1'], type='cmp')
                        ConfigPpt.fetch_data(shape.table.cell(x, 4), formed_data[region]['db_class'][item]['month_data_1'], formed_data[region]['db_class'][item]['month_data_0'], type='cmp')
                        
            # 페이지 복사
            if len(formed_data) > region_count:
                page_num = ConfigPpt.copy_slide(page_num)
    
    def elb(event, month):
        # 페이지 위치 / 데이터가 없을 경우 숨기기
        page_num = event['page_num']
        if 'data' not in event:
            ConfigPpt.hide_slide(page_num)
            return prs
        
        # 테이블 값 정리
        start_x = 1
        start_y = 0
        max_x = event['table']['LB_table']['max_x']
        max_y = event['table']['LB_table']['max_y']
        
        start_2nd_x = 1
        start_2nd_y = 1
        max_2nd_x = event['table']['TG_table']['max_x']
        max_2nd_y = event['table']['TG_table']['max_y']
        
        start_3rd_x = 1
        start_3rd_y = 1
        max_3rd_x = event['table']['T_table']['max_x']
        max_3rd_y = event['table']['T_table']['max_y']
        
        # region을 key로 데이터 재생성
        formed_data = ConfigData.event_format(event, region=True)
        
        # 슬라이드 개수 = region개수 
        for region_count, region in enumerate(list(formed_data), start=1):
            now_slide = prs.slides[page_num]
            
            # healthy, unhealthy 분리
            health_status = {}
            if 'healthy' in formed_data[region]['targetgroup']:
                health_status['healthy'] = formed_data[region]['targetgroup'].pop('healthy')
            if 'unhealthy' in formed_data[region]['targetgroup']:
                health_status['unhealthy'] = formed_data[region]['targetgroup'].pop('unhealthy')
                
            # target_total, group_total 분리
            target_group_total = {}
            if 'group_total' in formed_data[region]['targetgroup']:
                target_group_total['group_total'] = formed_data[region]['targetgroup'].pop('group_total')
            if 'target_total' in formed_data[region]['targetgroup']:
                target_group_total['target_total'] = formed_data[region]['targetgroup'].pop('target_total')
            
            # 슬라이드의 각 shape수정
            for shape in now_slide.shapes:
                if shape.name == 'Text Placeholder 1':
                    now_slide.shapes.element.remove(shape.element)
                    
                elif shape.name == 'region':
                    shape.text = region
                    
                elif shape.name == 'LB_chart':
                    chart_data = ChartData()
                    chart_data.categories = list([x.upper() for x in formed_data[region]['loadbalance']])
                    chart_data.add_series(month[0], tuple(x['month_data_3'] for x in formed_data[region]['loadbalance'].values()))
                    chart_data.add_series(month[1], tuple(x['month_data_2'] for x in formed_data[region]['loadbalance'].values()))
                    chart_data.add_series(month[2], tuple(x['month_data_1'] for x in formed_data[region]['loadbalance'].values()))
                    chart_data.add_series(month[3], tuple(x['month_data_0'] for x in formed_data[region]['loadbalance'].values()))
                    shape.chart.replace_data(chart_data)
                    
                elif shape.name == 'T_chart':
                    chart_data = ChartData()
                    chart_data.categories = month[::-1] # 역순
                    chart_data.add_series('Healthy', tuple(x for x in health_status['healthy'].values()))
                    chart_data.add_series('Unhealthy', tuple(x for x in health_status['unhealthy'].values()))
                    shape.chart.replace_data(chart_data)
                
                elif shape.name == 'LB_table':
                    # Initialize / 날짜 삽입
                    ConfigPpt.table_initialize(shape, start_x, start_y, max_x, max_y)
                    ConfigPpt.rename_date(shape, month, 0, 1)
                    
                    # 데이터 삽입
                    for x, item in enumerate(list(formed_data[region]['loadbalance']), start=start_x):
                        if x >= max_x:
                            break
                        ConfigPpt.fetch_data(shape.table.cell(x, 0), item.upper())
                        ConfigPpt.fetch_data(shape.table.cell(x, 1), formed_data[region]['loadbalance'][item]['month_data_3'])
                        ConfigPpt.fetch_data(shape.table.cell(x, 2), formed_data[region]['loadbalance'][item]['month_data_3'], formed_data[region]['loadbalance'][item]['month_data_2'], type='cmp')
                        ConfigPpt.fetch_data(shape.table.cell(x, 3), formed_data[region]['loadbalance'][item]['month_data_2'], formed_data[region]['loadbalance'][item]['month_data_1'], type='cmp')
                        ConfigPpt.fetch_data(shape.table.cell(x, 4), formed_data[region]['loadbalance'][item]['month_data_1'], formed_data[region]['loadbalance'][item]['month_data_0'], type='cmp')
                            
                elif shape.name == 'TG_table':
                    # Initialize / 날짜 삽입
                    ConfigPpt.table_initialize(shape, start_2nd_x, start_2nd_y, max_2nd_x, max_2nd_y)
                    ConfigPpt.rename_date(shape, month, 0, 1)
                    
                    # 데이터 삽입 (순서는 위에서 설정함)
                    for x, item in enumerate(list(target_group_total), start=start_2nd_x):
                        ConfigPpt.fetch_data(shape.table.cell(x, 1), target_group_total[item]['month_data_3'])
                        ConfigPpt.fetch_data(shape.table.cell(x, 2), target_group_total[item]['month_data_3'], target_group_total[item]['month_data_2'], type='cmp')
                        ConfigPpt.fetch_data(shape.table.cell(x, 3), target_group_total[item]['month_data_2'], target_group_total[item]['month_data_1'], type='cmp')
                        ConfigPpt.fetch_data(shape.table.cell(x, 4), target_group_total[item]['month_data_1'], target_group_total[item]['month_data_0'], type='cmp')
                
                elif shape.name == 'T_table':
                    # Initialize / 날짜 삽입
                    ConfigPpt.table_initialize(shape, start_3rd_x, start_3rd_y, max_3rd_x, max_3rd_y)
                    ConfigPpt.rename_date(shape, month, 0, 1)
                    
                    # 데이터 삽입 (순서때문에 for 사용 X)
                    for x, item in enumerate(list(health_status), start=start_3rd_x):
                        ConfigPpt.fetch_data(shape.table.cell(x, 1), health_status[item]['month_data_3'])
                        ConfigPpt.fetch_data(shape.table.cell(x, 2), health_status[item]['month_data_3'], health_status[item]['month_data_2'], type='cmp')
                        ConfigPpt.fetch_data(shape.table.cell(x, 3), health_status[item]['month_data_2'], health_status[item]['month_data_1'], type='cmp')
                        ConfigPpt.fetch_data(shape.table.cell(x, 4), health_status[item]['month_data_1'], health_status[item]['month_data_0'], type='cmp')
                        
            # 페이지 복사
            if len(formed_data) > region_count:
                page_num = ConfigPpt.copy_slide(page_num)
    
    def ebs(event, month):
        # 페이지 위치 / 데이터가 없을 경우 숨기기
        page_num = event['page_num']
        if 'data' not in event:
            ConfigPpt.hide_slide(page_num)
            return prs
        
        # 테이블 값 정리
        start_x = 2
        start_y = 2
        max_x = event['table']['EBS_table']['max_x']
        max_y = event['table']['EBS_table']['max_y']
        
        start_2nd_x = 1
        start_2nd_y = 0
        max_2nd_x = event['table']['EBS_type_table']['max_x']
        max_2nd_y = event['table']['EBS_type_table']['max_y']
        
        # {$region: {count: {...}}, size: {...}}, type: {...}}} 형태로 변형
        # volume not used 데이터 뽑아야함
        formed_data = {}
        for since_item in event['data']:
            if since_item not in formed_data:
                formed_data[since_item] = {}
            for region in event['data'][since_item]:
                if region not in formed_data[since_item]:
                    formed_data[since_item][region] = {}
                for target in event['data'][since_item][region]:
                    for type in event['data'][since_item][region][target]:
                        if type == 'type':
                            formed_data[since_item][region][type] = event['data'][since_item][region][target][type]
                            continue
                        for detail in event['data'][since_item][region][target][type]:
                            if detail not in formed_data[since_item][region]:
                                formed_data[since_item][region][detail] = {}
                            formed_data[since_item][region][detail][target+'-'+type] = event['data'][since_item][region][target][type][detail]
        
        # type 키 maxlen 초과시 other로 표기
        formed_data = ConfigData.event_format({'data': formed_data}, region=True)
        formed_data = ConfigData.data_maxlen_format(formed_data, region=True, key='type', maxlen=max_2nd_x-start_2nd_x)
        
        # sort formed_data
        sort_keys = ['snapshot-total', 'volume-total', 'volume-not_used']
        for region in formed_data:
            if 'count' in formed_data[region]:
                for sort_key in sort_keys:
                    if sort_key in formed_data[region]['count']:
                        formed_data[region]['count'][sort_key] = formed_data[region]['count'].pop(sort_key)
            if 'size' in formed_data[region]:
                for sort_key in sort_keys:
                    if sort_key in formed_data[region]['size']:
                        formed_data[region]['size'][sort_key] = formed_data[region]['size'].pop(sort_key)
                    
        # 슬라이드 개수 = region개수 
        for region_count, region in enumerate(list(formed_data), start=1):
            now_slide = prs.slides[page_num]
            # 슬라이드의 각 shape수정
            for shape in now_slide.shapes:
                if shape.name == 'Text Placeholder 1':
                    now_slide.shapes.element.remove(shape.element)
                    
                elif shape.name == 'region':
                    shape.text = region
                    
                elif shape.name == 'EBS_chart':
                    chart_data = ChartData()
                    chart_data.categories = list(formed_data[region]['size'])
                    chart_data.add_series(month[0], tuple(x['month_data_3'] for x in formed_data[region]['count'].values()))
                    chart_data.add_series(month[1], tuple(x['month_data_2'] for x in formed_data[region]['count'].values()))
                    chart_data.add_series(month[2], tuple(x['month_data_1'] for x in formed_data[region]['count'].values()))
                    chart_data.add_series(month[3], tuple(x['month_data_0'] for x in formed_data[region]['count'].values()))
                    shape.chart.replace_data(chart_data)
                    
                elif shape.name == 'EBS_type_chart':
                    chart_data = ChartData()
                    chart_data.categories = list(formed_data[region]['type'])
                    chart_data.add_series(month[3], tuple(x['month_data_0'] for x in formed_data[region]['type'].values()))
                    shape.chart.replace_data(chart_data)
                    
                elif shape.name == 'EBS_table':
                    # Initialize / 날짜 삽입
                    ConfigPpt.table_initialize(shape, start_x, start_y, max_x, max_y)
                    ConfigPpt.rename_date(shape, month, 1, 2)
                    # ConfigPpt.rename_date(shape, month, 1, 6)
                    
                    # 데이터 삽입
                    for x, item in enumerate(list(formed_data[region]['count']), start=start_x):
                        if x >= max_x:
                            break
                        ConfigPpt.fetch_data(shape.table.cell(x, 2), formed_data[region]['count'][item]['month_data_3'], zerofill='no')
                        ConfigPpt.fetch_data(shape.table.cell(x, 3), formed_data[region]['count'][item]['month_data_3'], formed_data[region]['count'][item]['month_data_2'], type='cmp', zerofill='no')
                        ConfigPpt.fetch_data(shape.table.cell(x, 4), formed_data[region]['count'][item]['month_data_2'], formed_data[region]['count'][item]['month_data_1'], type='cmp', zerofill='no')
                        ConfigPpt.fetch_data(shape.table.cell(x, 5), formed_data[region]['count'][item]['month_data_1'], formed_data[region]['count'][item]['month_data_0'], type='cmp', zerofill='no')
                        # ConfigPpt.fetch_data(shape.table.cell(x, 6), formed_data[region]['size'][item]['month_data_3'], zerofill='no')
                        # ConfigPpt.fetch_data(shape.table.cell(x, 7), formed_data[region]['size'][item]['month_data_2'], zerofill='no')
                        # ConfigPpt.fetch_data(shape.table.cell(x, 8), formed_data[region]['size'][item]['month_data_1'], zerofill='no')
                        # ConfigPpt.fetch_data(shape.table.cell(x, 9), formed_data[region]['size'][item]['month_data_0'], zerofill='no')
                        
                        # Volume Not Used
                        for y in range(max_y):
                            if x == 4:
                                ConfigPpt.text_red_bold(shape.table.cell(x,y))
                                
                elif shape.name == 'EBS_type_table':
                    # Initialize / 날짜 삽입
                    ConfigPpt.table_initialize(shape, start_2nd_x, start_2nd_y, max_2nd_x, max_2nd_y)
                    ConfigPpt.rename_date(shape, month, 0, 1)
                    
                    # 데이터 삽입
                    for x, item in enumerate(list(formed_data[region]['type']), start=start_2nd_x):
                        if x >= max_2nd_x:
                            break
                        ConfigPpt.fetch_data(shape.table.cell(x, 0), item)
                        ConfigPpt.fetch_data(shape.table.cell(x, 1), formed_data[region]['type'][item]['month_data_3'])
                        ConfigPpt.fetch_data(shape.table.cell(x, 2), formed_data[region]['type'][item]['month_data_3'], formed_data[region]['type'][item]['month_data_2'], type='cmp')
                        ConfigPpt.fetch_data(shape.table.cell(x, 3), formed_data[region]['type'][item]['month_data_2'], formed_data[region]['type'][item]['month_data_1'], type='cmp')
                        ConfigPpt.fetch_data(shape.table.cell(x, 4), formed_data[region]['type'][item]['month_data_1'], formed_data[region]['type'][item]['month_data_0'], type='cmp')
            
            # 페이지 복사
            if len(formed_data) > region_count:
                page_num = ConfigPpt.copy_slide(page_num)

    def ec2(event, month):
        # 페이지 위치 / 데이터가 없을 경우 숨기기
        page_num = event['page_num']
        if 'data' not in event:
            ConfigPpt.hide_slide(page_num)
            return prs
        
        # 테이블 값 정리
        start_x = 1
        start_y = 0
        max_x = event['table']['ec2_table']['max_x']
        max_y = event['table']['ec2_table']['max_y']
        
        start_2nd_x = 1
        start_2nd_y = 0
        max_2nd_x = event['table']['ec2_type_table']['max_x']
        max_2nd_y = event['table']['ec2_type_table']['max_y']
        
        # {instance: int} => {instance: {count: int}} 형태로 변형
        for since_item in event['data']:
            for region in event['data'][since_item]:
                temp_data = {}
                if 'instance' in event['data'][since_item][region]:
                    temp_data = {'count': event['data'][since_item][region]['instance']}
                event['data'][since_item][region].pop('instance')
                event['data'][since_item][region]['instance'] = temp_data
        
        # region을 key로 데이터 재생성
        formed_data = ConfigData.event_format(event, region=True)
        formed_data = ConfigData.data_maxlen_format(formed_data, region=True, key='instance_type', maxlen=max_2nd_x-start_2nd_x)
        
        # 슬라이드 개수 = region개수 
        for region_count, region in enumerate(list(formed_data), start=1):
            now_slide = prs.slides[page_num]
            # 슬라이드의 각 shape수정
            for shape in now_slide.shapes:
                if shape.name == 'Text Placeholder 1':
                    now_slide.shapes.element.remove(shape.element)
                    
                elif shape.name == 'region':
                    shape.text = region
                    
                elif shape.name == 'ec2_chart':
                    chart_data = ChartData()
                    chart_data.categories = list(formed_data[region]['instance_type'])
                    chart_data.add_series(month[0], tuple(x['month_data_3'] for x in formed_data[region]['instance_type'].values()))
                    chart_data.add_series(month[1], tuple(x['month_data_2'] for x in formed_data[region]['instance_type'].values()))
                    chart_data.add_series(month[2], tuple(x['month_data_1'] for x in formed_data[region]['instance_type'].values()))
                    chart_data.add_series(month[3], tuple(x['month_data_0'] for x in formed_data[region]['instance_type'].values()))
                    shape.chart.replace_data(chart_data)
                    
                elif shape.name == 'ec2_table':
                    # Initialize / 날짜 삽입
                    ConfigPpt.table_initialize(shape, start_x, start_y, max_x, max_y)
                    ConfigPpt.rename_date(shape, month, 0, 1)
                    
                    # 데이터 삽입
                    for x, item in enumerate(list(formed_data[region]['instance']), start=start_x):
                        if x >= max_x:
                            break
                        ConfigPpt.fetch_data(shape.table.cell(x, 1), formed_data[region]['instance'][item]['month_data_3'])
                        ConfigPpt.fetch_data(shape.table.cell(x, 2), formed_data[region]['instance'][item]['month_data_3'], formed_data[region]['instance'][item]['month_data_2'], type='cmp')
                        ConfigPpt.fetch_data(shape.table.cell(x, 3), formed_data[region]['instance'][item]['month_data_2'], formed_data[region]['instance'][item]['month_data_1'], type='cmp')
                        ConfigPpt.fetch_data(shape.table.cell(x, 4), formed_data[region]['instance'][item]['month_data_1'], formed_data[region]['instance'][item]['month_data_0'], type='cmp')
                        
                elif shape.name == 'ec2_type_table':
                    # Initialize / 날짜 삽입
                    ConfigPpt.table_initialize(shape, start_2nd_x, start_2nd_y, max_2nd_x, max_2nd_y)
                    ConfigPpt.rename_date(shape, month, 0, 1)
                    
                    # 데이터 삽입
                    for x, item in enumerate(list(formed_data[region]['instance_type']), start=start_2nd_x):
                        if x >= max_2nd_x:
                            break
                        ConfigPpt.fetch_data(shape.table.cell(x, 0), item)
                        ConfigPpt.fetch_data(shape.table.cell(x, 1), formed_data[region]['instance_type'][item]['month_data_3'])
                        ConfigPpt.fetch_data(shape.table.cell(x, 2), formed_data[region]['instance_type'][item]['month_data_3'], formed_data[region]['instance_type'][item]['month_data_2'], type='cmp')
                        ConfigPpt.fetch_data(shape.table.cell(x, 3), formed_data[region]['instance_type'][item]['month_data_2'], formed_data[region]['instance_type'][item]['month_data_1'], type='cmp')
                        ConfigPpt.fetch_data(shape.table.cell(x, 4), formed_data[region]['instance_type'][item]['month_data_1'], formed_data[region]['instance_type'][item]['month_data_0'], type='cmp')
            
            # 페이지 복사
            if len(formed_data) > region_count:
                page_num = ConfigPpt.copy_slide(page_num)
    
    def iam_user_list(event):
        # 페이지 위치 / 데이터가 없을 경우 숨기기
        page_num = event['page_num']
        if 'data' not in event:
            ConfigPpt.hide_slide(page_num)
            return prs
        else:
            if len(event['data']['month_data_0']) == 0:
                ConfigPpt.hide_slide(page_num)
                return prs
        
        # 테이블 값 정리
        start_x = 1
        start_y = 0
        max_x = event['table']['iam_userlist_table']['max_x']
        max_y = event['table']['iam_userlist_table']['max_y']
        repeat_counter = 0
        repeat_max_counter = 0
        
        # 2차원 배열로 변경 / 사용하지 않는 key 제거
        event['data'].pop('month_data_3', None)
        event['data'].pop('month_data_2', None)
        event['data'].pop('month_data_1', None)
        formed_data = []
        for num in range(math.ceil(len(event['data']['month_data_0'])/(max_x-start_x))):
            formed_data.append(event['data']['month_data_0'][num*(max_x-start_x) : (max_x-start_x)+num*(max_x-start_x)])
        repeat_max_counter  = len(formed_data)
        
        # 반복횟수 = len(event['data']['month_data_0']) = 2차원 배열 row수
        for array_counter in range(repeat_max_counter):
            now_slide = prs.slides[page_num]
            repeat_counter += 1
            for shape in now_slide.shapes:
                if shape.name == 'iam_userlist_table':
                    ConfigPpt.table_initialize(shape, start_x, start_y, max_x, max_y)
                    for x, item in enumerate(formed_data[array_counter], start=start_x):
                        shape.table.cell(x,0).text = str(item['user'])
                        shape.table.cell(x,1).text = str(item['password_enabled'])
                        shape.table.cell(x,2).text = str(item['password_last_used'])
                        shape.table.cell(x,3).text = str(item['password_last_changed'])
                        shape.table.cell(x,4).text = str(item['mfa_active'])
                        shape.table.cell(x,5).text = str(item['access_key_count'])+' ea'
                        shape.table.cell(x,6).text = str(item['access_key_used_date'])
                        
                        # 글씨 속성 변경
                        for y in range(max_y):
                            ConfigPpt.table_font(shape.table.cell(x,y))
                            if y == 0:
                                ConfigPpt.text_bold(shape.table.cell(x,y))
                            if shape.table.cell(x,y).text == '미사용':
                                ConfigPpt.text_red_bold(shape.table.cell(x,y))
                            if y == 4 and item['password_enabled'] == '활성화' and item['mfa_active'] == '비활성화':
                                ConfigPpt.text_red_bold(shape.table.cell(x,y))
                            if y == 5 and item['access_key_count'] > 1:
                                ConfigPpt.text_red_bold(shape.table.cell(x,y))
                            if y != 6 and '일 전' in shape.table.cell(x,y).text:
                                int_date = int(shape.table.cell(x,y).text.split('일')[0])
                                if int_date > 90:
                                    ConfigPpt.text_red_bold(shape.table.cell(x,y))
                            if y == 6 and '일 전' in shape.table.cell(x,y).text:
                                int_date = int(shape.table.cell(x,y).text.split('일')[0])
                                if int_date > 180:
                                    ConfigPpt.text_red_bold(shape.table.cell(x,y))
            
            # 마지막은 복사 하지 않기        
            if repeat_counter != repeat_max_counter:
                page_num = ConfigPpt.copy_slide(page_num)
            else:
                return prs
    
    
def lambda_handler(event, context):
    
    # 로깅
    logger.info('{}'.format(json.dumps(event)))
    logger.info('{}-{}({}) Create PPT Start.'.format(
        event['customer_name'], event['customer_project'], event['account_id']))
        
    # 날짜 지정
    # run_date = date.today()
    # run_date = '2022-03-30'
    run_date = event['run_date']
    report_cycle = event['report_cycle']
    report_date = ConfigData.date_format(run_date, report_cycle)
    
    # s3에서 PPT 파일 가져오기
    global prs 
    prs = ConfigPpt.presentation()
    
    # prs의 각 내용 가져오기
    slide_idx = ConfigPpt.index_ppt(prs)
    
    # this_month / last_month / before_last_month => agenda_list
    slide_idx = ConfigData.add_event_data(event, slide_idx)
    
    # FixPage 
    FixPage.report_intro_agenda(slide_idx['agenda'])
    FixPage.main(slide_idx['main'], report_date['main'], report_cycle)
    FixPage.information(slide_idx['information'], report_date['info'])
    FixPage.iam_summary(slide_idx['iam_summary'], report_date['table'])
    FixPage.work_summary(slide_idx['work_summary'])
    
    # # ElasticPage
    ElasticPage.trusted_advisor(slide_idx['trusted_advisor'])
    ElasticPage.ri_usage(slide_idx['ri_usage'])
    ElasticPage.ri_info(slide_idx['ri_info'])
    ElasticPage.backup(slide_idx['backup'])
    ElasticPage.work_list(slide_idx['work_list'])
    ElasticPage.other(slide_idx['other'], report_date['table'])
    ElasticPage.route53(slide_idx['route53'])
    ElasticPage.cloudfront(slide_idx['cloud_front'], report_date['table'])
    ElasticPage.cache(slide_idx['cache'], report_date['table'])
    ElasticPage.rds_single(slide_idx['rds_single'], report_date['table'])
    ElasticPage.rds_cluster(slide_idx['rds_cluster'], report_date['table'])
    ElasticPage.elb(slide_idx['elb'], report_date['table'])
    ElasticPage.ebs(slide_idx['ebs'], report_date['table'])
    ElasticPage.ec2(slide_idx['ec2'], report_date['table'])
    ElasticPage.iam_user_list(slide_idx['iam_list'])
            
    logger.info('{}-{}({}) Create PPT Finish! :> Total {} Pages'.format(
        event['customer_name'], event['customer_project'], 
        event['account_id'], len(prs.slides)))
        
    # S3_Upload
    try:
        with io.BytesIO() as out:
            prs.save(out)
            out.seek(0)
            today = datetime.now()
            """ Report Cycle 영어로 할 경우 하단 Value 변경 (ex - 월간 -> monthly) / 파일 저장할 위치 지정"""
            if report_cycle == '월간':
                response = boto3.client('s3').upload_fileobj(out, 'FIX_INPUT: Bucket명', 'monthly/'+str(date.today())+'/'+event['customer_name']+'/'+event['customer_name']+'-'+event['customer_project']+'_'+report_date['file']+'_'+event['report_cycle']+'보고서.pptx')
            elif report_cycle == '분기':
                response = boto3.client('s3').upload_fileobj(out, 'FIX_INPUT: Bucket명', 'quarter/'+str(date.today())+'/'+event['customer_name']+'/'+event['customer_name']+'-'+event['customer_project']+'_'+report_date['file']+'_'+event['report_cycle']+'보고서.pptx')
            else:
                print('보고서 주기 확인 불가')
    except Exception as e:
        logger.error('[S3_Upload] {}'.format(e))
    else:
        logger.info('[{}-{}({})] S3 Upload Finish ~'.format(
            event['customer_name'], event['customer_project'], 
            event['account_id']))
